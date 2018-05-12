from pptx import Presentation
from pptx.util import Inches

# 获取豆瓣电影信息
def getInfo():
	movies_info = []

	with open('./douban_movie_top250.txt') as f:
		for line in f.readlines():
			line_list = line.split("\'")
			one_movie_info = {}
			one_movie_info['index'] = line_list[1]
			one_movie_info['title'] = line_list[3]
			one_movie_info['score'] = line_list[5]

			try:
				one_movie_info['desc'] = line_list[7]
			except:
				one_movie_info['desc'] = ''

			one_movie_info['image_path'] = "./Top250_movie_images/"+ str(line_list[1]) + '_' + line_list[3] + ".jpg"

			movies_info.append(one_movie_info)

	return movies_info


# 创建ppt
def createPpt(movies_info):
	prs = Presentation('model.pptx')
	for movie_info in movies_info:
		# 获取模板个数
		templateStyleNum = len(prs.slide_layouts)
		# 按照第一个模板创建 一张幻灯片
		oneSlide = prs.slides.add_slide(prs.slide_layouts[0])
		# 获取模板可填充的所有位置
		body_shapes = oneSlide.shapes.placeholders
		for index, body_shape in enumerate(body_shapes):
			if index == 0:
				body_shape.text = movie_info['index']+movie_info['title']
			elif index == 1:
				img_path = movie_info['image_path']
				body_shape.insert_picture(img_path)
			elif index == 2:
				body_shape.text = movie_info['desc']
			elif index == 3:
				body_shape.text = movie_info['score']
	# 对ppt的修改  
	prs.save('豆瓣Top250推荐.pptx')



def main():
	# 获取豆瓣电影信息
	movies_info = getInfo()
	createPpt(movies_info)


if __name__ == '__main__':
	main()