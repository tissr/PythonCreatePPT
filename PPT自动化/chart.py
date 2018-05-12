# encoding: utf-8
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

# 创建幻灯片
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 定义图表数据
chart_data = ChartData()
chart_data.categories = ['肖申克的救赎', '霸王别姬', '这个杀手不太冷', '阿甘正传', '美丽人生', '千与千寻', '泰坦尼克号', '辛德勒的名单', '盗梦空间', '机器人总动员', '海上钢琴师', '三傻大闹宝莱坞', '忠犬八公的故事', '放牛班的春天', '大话西游之大圣娶亲', '楚门的世界', '教父', '龙猫', '熔炉', '乱世佳人']
chart_data.add_series('豆瓣电影', (9.6, 9.5, 9.4, 9.4, 9.5, 9.2, 9.2, 9.4, 9.3, 9.3, 9.2, 9.2, 9.2, 9.2, 9.2, 9.1, 9.2, 9.1, 9.2, 9.2))

# 将图表添加到幻灯片
x, y, cx, cy = Inches(0), Inches(0), Inches(10), Inches(8)
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

prs.save('豆瓣 Top20 评分图.pptx')